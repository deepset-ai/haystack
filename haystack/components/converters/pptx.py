# SPDX-FileCopyrightText: 2022-present deepset GmbH <info@deepset.ai>
#
# SPDX-License-Identifier: Apache-2.0

import io
import os
from pathlib import Path
from typing import Any, Literal

from haystack import Document, component, default_to_dict, logging
from haystack.components.converters.utils import get_bytestream_from_source, normalize_metadata
from haystack.dataclasses import ByteStream
from haystack.lazy_imports import LazyImport

with LazyImport("Run 'pip install python-pptx'") as pptx_import:
    from pptx import Presentation
    from pptx.text.text import _Paragraph


logger = logging.getLogger(__name__)


@component
class PPTXToDocument:
    """
    Converts PPTX files to Documents.

    Usage example:
    ```python
    from haystack.components.converters.pptx import PPTXToDocument

    converter = PPTXToDocument()
    results = converter.run(sources=["sample.pptx"], meta={"date_added": datetime.now().isoformat()})
    documents = results["documents"]
    print(documents[0].content)
    # 'This is the text from the PPTX file.'
    ```
    """

    def __init__(self, store_full_path: bool = False, link_format: Literal["markdown", "plain", "none"] = "none"):
        """
        Create a PPTXToDocument component.

        :param store_full_path:
            If True, the full path of the file is stored in the metadata of the document.
            If False, only the file name is stored.
        :param link_format: The format for link output. Possible options:
            - `"markdown"`: `[text](url)`
            - `"plain"`: `text (url)`
            - `"none"`: Only the text is extracted, link addresses are ignored.
        """
        pptx_import.check()
        if link_format not in ("markdown", "plain", "none"):
            msg = f"Unknown link format '{link_format}'. Supported formats are: 'markdown', 'plain', 'none'"
            raise ValueError(msg)
        self.link_format = link_format
        self.store_full_path = store_full_path

    def to_dict(self) -> dict[str, Any]:
        """
        Serializes the component to a dictionary.

        :returns:
            Dictionary with serialized data.
        """
        return default_to_dict(self, link_format=self.link_format, store_full_path=self.store_full_path)

    def _convert(self, file_content: io.BytesIO) -> str:
        """
        Converts the PPTX file to text.
        """
        pptx_presentation = Presentation(file_content)
        text_all_slides = []
        for slide in pptx_presentation.slides:
            text_on_slide = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paragraphs = []
                    for paragraph in shape.text_frame.paragraphs:
                        paragraphs.append(self._process_paragraph(paragraph))
                    text_on_slide.append("\n".join(paragraphs))
                elif hasattr(shape, "text"):
                    text_on_slide.append(shape.text)
            text_all_slides.append("\n".join(text_on_slide))
        return "\f".join(text_all_slides)

    def _process_paragraph(self, paragraph: "_Paragraph") -> str:
        """
        Processes a paragraph and formats hyperlinks according to the specified link format.

        :param paragraph: The PPTX paragraph to process.
        :returns: A string with links formatted according to the specified format.
        """
        if self.link_format == "none":
            return paragraph.text
        parts = []
        for run in paragraph.runs:
            if run.hyperlink and run.hyperlink.address:
                if self.link_format == "markdown":
                    parts.append(f"[{run.text}]({run.hyperlink.address})")
                else:
                    parts.append(f"{run.text} ({run.hyperlink.address})")
            else:
                parts.append(run.text)
        return "".join(parts)

    @component.output_types(documents=list[Document])
    def run(self, sources: list[str | Path | ByteStream], meta: dict[str, Any] | list[dict[str, Any]] | None = None):
        """
        Converts PPTX files to Documents.

        :param sources:
            List of file paths or ByteStream objects.
        :param meta:
            Optional metadata to attach to the Documents.
            This value can be either a list of dictionaries or a single dictionary.
            If it's a single dictionary, its content is added to the metadata of all produced Documents.
            If it's a list, the length of the list must match the number of sources, because the two lists will
            be zipped.
            If `sources` contains ByteStream objects, their `meta` will be added to the output Documents.

        :returns:
            A dictionary with the following keys:
            - `documents`: Created Documents
        """
        documents = []
        meta_list = normalize_metadata(meta, sources_count=len(sources))

        for source, metadata in zip(sources, meta_list, strict=False):
            try:
                bytestream = get_bytestream_from_source(source)
            except Exception as e:
                logger.warning("Could not read {source}. Skipping it. Error: {error}", source=source, error=e)
                continue
            try:
                text = self._convert(io.BytesIO(bytestream.data))
            except Exception as e:
                logger.warning(
                    "Could not read {source} and convert it to Document, skipping. {error}", source=source, error=e
                )
                continue

            merged_metadata = {**bytestream.meta, **metadata}

            if not self.store_full_path and (file_path := bytestream.meta.get("file_path")):
                merged_metadata["file_path"] = os.path.basename(file_path)
            documents.append(Document(content=text, meta=merged_metadata))

        return {"documents": documents}
